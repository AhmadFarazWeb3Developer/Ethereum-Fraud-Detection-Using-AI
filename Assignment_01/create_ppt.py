from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1]
content.text = "• Ethereum blockchain faces significant fraud risks from malicious actors\n• Need to detect fraudulent addresses based on transactional behavior\n• Data collection and cleaning is the foundation for fraud detection models\n• Addresses with abnormal transaction patterns indicate potential fraud"

# Slide 2: Datasets
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Datasets (Sources, Descriptions etc.)"
content = slide.placeholders[1]
content.text = "• Dataset: Ethereum_Fraud_Detection.csv\n• Source: Blockchain transaction data (specific source not detailed)\n• Description: Contains transactional features for Ethereum addresses\n• Features include: transaction counts, ETH values, ERC20 token interactions\n• Target variable: FLAG (1 = fraudulent, 0 = legitimate)\n• Size: Approximately 9,000+ records"

# Slide 3: Analytical Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Analytical Approach"
content = slide.placeholders[1]
content.text = "• Data Collection: Load and examine the raw dataset\n• Data Exploration: Check data types, shapes, and basic statistics\n• Data Cleaning: Handle missing values, duplicates, and inconsistencies\n• Validation: Ensure data quality for subsequent analysis\n• Storage: Save cleaned dataset for future assignments"

# Slide 4: Your Project Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Project Methodology"
content = slide.placeholders[1]
content.text = "• Step 1: Import necessary libraries (pandas)\n• Step 2: Load the dataset using pd.read_csv()\n• Step 3: Examine dataset structure (shape, head, dtypes)\n• Step 4: Check for missing values in numeric and object columns\n• Step 5: Identify and handle duplicate records\n• Step 6: Clean null, empty, and invalid values\n• Step 7: Save the cleaned dataset"

# Slide 5: Preprocessing Techniques
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Preprocessing Techniques"
content = slide.placeholders[1]
content.text = "• Missing Value Handling: Fill numeric NaNs with 0, object NaNs with '0'\n• String Cleaning: Strip whitespace and replace empty strings with '0'\n• Duplicate Removal: Check for duplicate addresses (found 25 duplicates)\n• Data Type Consistency: Ensure proper data types for analysis\n• Column Standardization: Clean column names for consistency"

# Slide 6: Models
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Models"
content = slide.placeholders[1]
content.text = "• No machine learning models developed in this assignment\n• Focus was on data preparation\n• Models will be built in subsequent assignments using the cleaned data"

# Slide 7: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Results"
content = slide.placeholders[1]
content.text = "• Dataset Shape: 9841 rows × 49 columns\n• No duplicate records found across all columns\n• 25 duplicate addresses identified\n• Missing values handled appropriately\n• Cleaned dataset saved as 'Cleaned_Ethereum_Fraud_Detection.csv'\n• Data ready for exploratory data analysis"

# Slide 8: Your Findings and Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Findings and Conclusions"
content = slide.placeholders[1]
content.text = "• Dataset contains comprehensive transactional features\n• Minimal data quality issues found\n• FLAG column has no missing values, indicating reliable target variable\n• ERC20 token interactions show some missing data patterns\n• Data cleaning preserved all important information\n• Foundation established for fraud detection analysis"

# Slide 9: Deployment Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Deployment Strategy"
content = slide.placeholders[1]
content.text = "• No deployment in this assignment\n• Cleaned data will be used in Assignment 02 for EDA\n• Data quality ensures reliable model training in future assignments"

# Slide 10: Feedback and its Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Feedback and its Analysis"
content = slide.placeholders[1]
content.text = "• No user feedback collected in this data collection phase\n• Feedback will be gathered in later assignments during deployment"

# Slide 11: Future Goals
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Goals"
content = slide.placeholders[1]
content.text = "• Perform Exploratory Data Analysis (EDA) in Assignment 02\n• Identify patterns and correlations in the data\n• Develop machine learning models for fraud detection\n• Deploy the final model for real-world use\n• Collect user feedback and iterate on improvements"

# Save the presentation
prs.save('Assignment_01_Presentation.pptx')
print("Assignment_01_Presentation.pptx created successfully")