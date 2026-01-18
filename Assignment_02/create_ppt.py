from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1]
content.text = "• Build upon cleaned data from Assignment 01\n• Perform Exploratory Data Analysis (EDA) to understand data patterns\n• Identify key features that distinguish fraudulent from legitimate addresses\n• Establish methodological foundation for fraud detection modeling\n• Uncover insights about Ethereum transaction behaviors"

# Slide 2: Datasets
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Datasets (Sources, Descriptions etc.)"
content = slide.placeholders[1]
content.text = "• Dataset: Cleaned_Ethereum_Fraud_Detection.csv (from Assignment 01)\n• Description: Preprocessed Ethereum transaction data\n• Features: 49 columns including transaction counts, ETH values, ERC20 metrics\n• Target: FLAG column (binary: 0=legitimate, 1=fraudulent)\n• Size: 9841 records\n• Data types: Mix of numeric (int64, float64) and categorical features"

# Slide 3: Analytical Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Analytical Approach"
content = slide.placeholders[1]
content.text = "• Dataset Overview: Examine structure, shape, and basic statistics\n• Descriptive Statistics: Analyze central tendencies and distributions\n• Target Analysis: Study class distribution and imbalance\n• Univariate Analysis: Examine individual feature distributions\n• Bivariate Analysis: Explore relationships between features and target\n• Correlation Analysis: Identify feature interdependencies\n• Outlier Detection: Use boxplots to identify anomalous values"

# Slide 4: Your Project Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Project Methodology"
content = slide.placeholders[1]
content.text = "• Load cleaned dataset using pandas\n• Perform dataset overview (shape, info, head)\n• Generate descriptive statistics with describe()\n• Analyze target variable distribution\n• Create histograms for feature distributions"

# Slide 4.5: Your Project Methodology (continued)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Project Methodology (continued)"
content = slide.placeholders[1]
content.text = "• Generate boxplots for outlier detection\n• Compute correlation matrix with heatmap visualization\n• Examine relationships with scatter plots\n• Document key findings and insights"

# Slide 6: Preprocessing Techniques
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Preprocessing Techniques"
content = slide.placeholders[1]
content.text = "• Data already cleaned in Assignment 01\n• Column name standardization applied\n• Missing values handled (filled with 0 or '0')\n• Duplicate records removed\n• Data types verified for analysis\n• Ready for EDA without additional preprocessing"

# Slide 7: Models
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Models"
content = slide.placeholders[1]
content.text = "• No predictive models developed in this EDA phase\n• Focus on understanding data patterns and distributions\n• Insights will inform model selection in Assignment 03\n• Statistical analysis rather than machine learning"

# Slide 8: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Results"
content = slide.placeholders[1]
content.text = "• Dataset: 9841 rows × 49 columns\n• Class Distribution: Significant imbalance (mostly legitimate addresses)\n• Feature Distributions: Right-skewed for transaction values\n• Outliers Detected: Many features show extreme values\n• Strong correlations between transaction-related features\n• Fraudulent addresses show distinct patterns in ERC20 activities\n• Key relationships identified between total transactions and fraud likelihood"

# Slide 9: Your Findings and Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Findings and Conclusions"
content = slide.placeholders[1]
content.text = "• Fraudulent addresses exhibit abnormal transaction behaviors\n• ERC20 token interactions are strong fraud indicators\n• Class imbalance requires special handling in modeling\n• Transaction frequency and value patterns distinguish fraud\n• Dataset is suitable for supervised classification\n• EDA confirms behavioral patterns for fraud detection"

# Slide 10: Deployment Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Deployment Strategy"
content = slide.placeholders[1]
content.text = "• No deployment in this exploratory phase\n• Findings will guide model development in Assignment 03\n• Insights inform feature selection and preprocessing strategies"

# Slide 11: Feedback and its Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Feedback and its Analysis"
content = slide.placeholders[1]
content.text = "• No user feedback collected during EDA phase\n• Analysis is internal to data understanding\n• Feedback will be gathered during model deployment"

# Slide 12: Future Goals
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Goals"
content = slide.placeholders[1]
content.text = "• Implement machine learning models in Assignment 03\n• Address class imbalance with appropriate techniques\n• Develop Logistic Regression and Random Forest models\n• Evaluate models using proper metrics (precision, recall, F1, ROC-AUC)\n• Deploy trained models for real-world fraud detection\n• Collect user feedback and iterate on improvements"

# Save the presentation
prs.save('Assignment_02_Presentation.pptx')
print("Assignment_02_Presentation.pptx created successfully")