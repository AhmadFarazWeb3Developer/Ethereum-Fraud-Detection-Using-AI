from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1]
content.text = "• Transition fraud detection models from development to production\n• Deploy machine learning models for real-world Ethereum fraud detection\n• Create user-friendly interface for fraud risk assessment\n• Collect and analyze user feedback for iterative improvements\n• Demonstrate complete data science project lifecycle"

# Slide 2: Datasets
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Datasets (Sources, Descriptions etc.)"
content = slide.placeholders[1]
content.text = "• Primary Dataset: Cleaned_Ethereum_Fraud_Detection.csv\n• Model Artifacts: random_forest_model.pkl, logistic_regression_model.pkl, scaler.pkl\n• Feedback Dataset: Suggestions_Dataset.csv (collected post-deployment)\n• Features: Transaction counts, ETH values, ERC20 metrics\n• Target: Binary fraud classification\n• Models trained in Assignment 03, deployed in Assignment 04"

# Slide 3: Analytical Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Analytical Approach"
content = slide.placeholders[1]
content.text = "• Model Deployment: Load pre-trained models from Assignment 03\n• Web Application Development: Streamlit for user interface\n• User Input Processing: Accept transaction features via web form\n• Prediction Generation: Real-time fraud risk assessment\n• Model Comparison: Display results from both RF and LR models\n• Feedback Collection: Google Forms for user suggestions\n• Feedback Analysis: Thematic analysis of improvement suggestions"

# Slide 4: Your Project Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Project Methodology"
content = slide.placeholders[1]
content.text = "• Load pre-trained models using pickle\n• Develop Streamlit web application\n• Implement user input interface for transaction features\n• Process inputs and generate predictions\n• Display results with probability scores and risk visualization\n• Compare Random Forest and Logistic Regression outputs\n• Collect user feedback via structured forms\n• Analyze feedback themes and prioritize improvements"

# Slide 5: Preprocessing Techniques
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Preprocessing Techniques"
content = slide.placeholders[1]
content.text = "• Model Loading: Deserialize pickle files from Assignment 03\n• Feature Consistency: Maintain same column order as training\n• Input Validation: Ensure numeric inputs for transaction features\n• Scaling Application: Use saved scaler for Logistic Regression\n• Data Formatting: Convert user inputs to numpy arrays\n• Error Handling: Graceful handling of missing models or invalid inputs"

# Slide 6: Models
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Models"
content = slide.placeholders[1]
content.text = "• Primary Model: Random Forest Classifier (200 estimators)\n  - Superior performance (ROC-AUC: 0.94)\n  - Better recall for fraud detection\n  - No feature scaling required\n• Baseline Model: Logistic Regression\n  - Interpretable results\n  - Requires StandardScaler\n  - Used for model comparison\n• Selection Criteria: Performance, recall, and deployment simplicity"

# Slide 7: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Results"
content = slide.placeholders[1]
content.text = "• Successful Streamlit deployment\n• Web interface accepts 4 key features:\n  - Total Transactions\n  - Total ETH Received\n  - Average Transaction Value\n  - Number of ERC20 Tokens\n• Real-time predictions with probability scores\n• Risk visualization with progress bars\n• Model comparison table (RF vs LR)\n• Application ready for local execution"

# Slide 8: Your Findings and Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Findings and Conclusions"
content = slide.placeholders[1]
content.text = "• Successful transition from model development to deployment\n• Streamlit provides effective platform for ML deployment\n• User interface enables practical fraud risk assessment\n• Model performance maintained in production environment\n• Feedback collection reveals usability and feature improvement opportunities\n• Complete data science lifecycle demonstrated\n• Foundation established for iterative model improvements"

# Slide 9: Deployment Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Deployment Strategy"
content = slide.placeholders[1]
content.text = "• Local Streamlit Application\n• Command: streamlit run app.py\n• User Interface: Web-based form for input\n• Real-time Processing: Immediate prediction results\n• Model Serialization: Pickle format for model persistence\n• Scalability: Suitable for academic and small-scale production use\n• Future: Could extend to cloud deployment (Heroku, AWS)"

# Slide 10: Feedback and its Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Feedback and its Analysis"
content = slide.placeholders[1]
content.text = "• Feedback collected from 15+ users via Google Forms\n• Key Themes:\n  - Usability improvements requested\n  - Additional input features suggested\n  - UI clarity enhancements needed\n  - Probability-based risk scoring valued\n• Analysis: Prioritized accuracy and usability improvements\n• Common Suggestions: More features, better explanations, enhanced UI"

# Slide 11: Future Goals
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Goals"
content = slide.placeholders[1]
content.text = "• Version 2.0 Development:\n  - Add more input features for better accuracy\n  - Improve UI with tooltips and better explanations\n  - Implement advanced risk scoring algorithms\n  - Add batch processing capabilities\n• Model Improvements:\n  - Regular retraining with new data\n  - Hyperparameter optimization\n  - Ensemble model development\n• Deployment Enhancements:\n  - Cloud deployment (Heroku, Streamlit Cloud)\n  - API development for integration\n  - Mobile application development"

# Save the presentation
prs.save('Assignment_04_Presentation.pptx')
print("Assignment_04_Presentation.pptx created successfully")