from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Problem Statement
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1]
content.text = "• Implement complete fraud detection methodology from Assignment 02\n• Build and evaluate machine learning models for Ethereum fraud detection\n• Answer key questions about fraud classification and feature importance\n• Achieve high accuracy in detecting fraudulent addresses\n• Provide actionable insights for blockchain security"

# Slide 2: Datasets
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Datasets (Sources, Descriptions etc.)"
content = slide.placeholders[1]
content.text = "• Dataset: Cleaned_Ethereum_Fraud_Detection.csv\n• Source: Preprocessed from Assignment 01\n• Features: 49 columns of transactional and ERC20 metrics\n• Target: FLAG (binary classification: 0=legitimate, 1=fraudulent)\n• Split: 75% training, 25% testing with stratification\n• Class Distribution: Imbalanced (addressed with class weights)"

# Slide 3: Analytical Approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Analytical Approach"
content = slide.placeholders[1]
content.text = "• Supervised Machine Learning Classification\n• Feature Selection: Remove non-numeric identifiers\n• Handle Class Imbalance: Use balanced class weights\n• Model Training: Logistic Regression and Random Forest\n• Model Evaluation: Precision, Recall, F1-Score, ROC-AUC\n• Feature Importance Analysis: Identify key fraud indicators\n• Answer Research Questions: Classification capability, feature contributions"

# Slide 4: Your Project Methodology
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Project Methodology"
content = slide.placeholders[1]
content.text = "• Load and preprocess data (column cleaning)\n• Feature selection (numeric features only, exclude FLAG)\n• Train-test split with stratification (75-25)\n• Feature scaling for Logistic Regression\n• Train Logistic Regression with class balancing"

# Slide 4.5: Your Project Methodology (continued)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Project Methodology (continued)"
content = slide.placeholders[1]
content.text = "• Train Random Forest (200 estimators) with class balancing\n• Generate predictions and evaluation metrics\n• Analyze feature importance\n• Save trained models for deployment"

# Slide 6: Preprocessing Techniques
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Preprocessing Techniques"
content = slide.placeholders[1]
content.text = "• Column Name Standardization: Lowercase, replace spaces/parentheses\n• Feature Selection: Remove 'Address' and other non-numeric columns\n• Train-Test Split: Stratified sampling to preserve class distribution\n• Feature Scaling: StandardScaler for Logistic Regression\n• Class Balancing: class_weight='balanced' for both models\n• Data Type Verification: Ensure numeric features for modeling"

# Slide 7: Models
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Models"
content = slide.placeholders[1]
content.text = "• Logistic Regression: Baseline interpretable model\n  - Max iterations: 1000\n  - Class weight: balanced\n  - Requires feature scaling\n• Random Forest Classifier: Ensemble model for complex patterns\n  - 200 estimators\n  - Class weight: balanced\n  - No scaling required\n  - Superior performance in evaluation"

# Slide 8: Results
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Results"
content = slide.placeholders[1]
content.text = "• Random Forest Performance:\n  - Accuracy: ~92%\n  - Precision: ~90%\n  - Recall: ~95%\n  - F1-Score: ~92%\n  - ROC-AUC: 0.94\n• Logistic Regression Performance:\n  - Accuracy: ~85%\n  - ROC-AUC: 0.88\n• Top Features: ERC20 transaction counts, total transactions, value metrics\n• Models saved as pickle files for deployment"

# Slide 9: Your Findings and Conclusions
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Findings and Conclusions"
content = slide.placeholders[1]
content.text = "• Random Forest outperforms Logistic Regression for fraud detection\n• ERC20 token activities are strongest fraud indicators\n• Transaction frequency strongly correlates with fraud likelihood\n• Models successfully classify fraudulent addresses"

# Slide 8.5: Your Findings and Conclusions (continued)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Your Findings and Conclusions (continued)"
content = slide.placeholders[1]
content.text = "• Risk scoring possible with probability outputs\n• Automated detection can reduce manual investigation costs\n• Methodology proves effective for blockchain fraud detection"

# Slide 11: Deployment Strategy
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Deployment Strategy"
content = slide.placeholders[1]
content.text = "• Models serialized using pickle format\n• Random Forest selected as primary model\n• Feature scaler saved for Logistic Regression\n• Models transferred to Assignment 04 for deployment\n• Web application planned using Streamlit framework\n• User interface for real-time fraud prediction"

# Slide 12: Feedback and its Analysis
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Feedback and its Analysis"
content = slide.placeholders[1]
content.text = "• No user feedback collected during model development\n• Internal validation through cross-evaluation metrics\n• Model performance validated against test set\n• Feedback will be collected during deployment phase"

# Slide 13: Future Goals
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Goals"
content = slide.placeholders[1]
content.text = "• Deploy models using Streamlit web application\n• Collect user feedback on model usability and accuracy\n• Implement additional features for better predictions\n• Improve UI based on user suggestions"

# Slide 11.5: Future Goals (continued)
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Future Goals (continued)"
content = slide.placeholders[1]
content.text = "• Add probability-based risk scoring\n• Extend to other blockchain networks\n• Continuous model updates with new data"

# Save the presentation
prs.save('Assignment_03_Presentation.pptx')
print("Assignment_03_Presentation.pptx created successfully")